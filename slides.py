import os
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

PDF_FOLDER = "pdfs"
OUTPUT_FOLDER = "output_ppt"

os.makedirs(OUTPUT_FOLDER, exist_ok=True)

for file in os.listdir(PDF_FOLDER):
    if not file.lower().endswith(".pdf"):
        continue

    pdf_path = os.path.join(PDF_FOLDER, file)
    ppt_name = file.replace(".pdf", "_locked.pptx")
    ppt_path = os.path.join(OUTPUT_FOLDER, ppt_name)

    print(f"Processing: {file}")

    images = convert_from_path(pdf_path, dpi=300)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_slide = prs.slide_layouts[6]

    for i, img in enumerate(images):
        slide = prs.slides.add_slide(blank_slide)
        img_path = f"temp_{i}.png"
        img.save(img_path, "PNG")

        slide.shapes.add_picture(
            img_path,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )

        os.remove(img_path)

    prs.save(ppt_path)
    print(f"Created: {ppt_path}")

print("✅ All PDFs converted successfully.")
